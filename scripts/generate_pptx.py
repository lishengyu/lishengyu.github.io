"""
Generate Design Proposal PPTX for Personal Website Project
Uses python-pptx library
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Color palette (Minimal & Modern)
COLOR_PRIMARY   = RGBColor(0x1a, 0x1a, 0x2e)  # Dark navy
COLOR_ACCENT    = RGBColor(0x4f, 0x46, 0xe5)  # Indigo
COLOR_ACCENT2   = RGBColor(0x06, 0xb6, 0xd4)  # Cyan
COLOR_BG        = RGBColor(0xf8, 0xf9, 0xfa)  # Light gray
COLOR_WHITE     = RGBColor(0xff, 0xff, 0xff)
COLOR_TEXT      = RGBColor(0x33, 0x33, 0x33)
COLOR_TEXT_LIGHT= RGBColor(0x66, 0x66, 0x66)
COLOR_SURFACE   = RGBColor(0xff, 0xff, 0xff)
COLOR_BORDER    = RGBColor(0xe2, 0xe8, 0xf0)

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)

# ---- Helper Functions ----

def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, left, top, width, height,
                font_size=14, bold=False, color=COLOR_TEXT,
                align=PP_ALIGN.LEFT, font_name='Arial'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

# ================================================================
# Slide 1: Title Slide
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
set_slide_bg(slide, COLOR_PRIMARY)

# Decorative accent bar (top)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(3.5), Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = COLOR_ACCENT2
bar.line.fill.background()

# Title
add_textbox(slide, '个人网站设计方案', 1, 1.2, 8, 1,
            font_size=36, bold=True, color=COLOR_WHITE, align=PP_ALIGN.LEFT)

# Subtitle
add_textbox(slide, 'GitHub Pages · 技术博客 · 简洁现代风格', 1, 2.3, 8, 0.6,
            font_size=18, bold=False, color=COLOR_ACCENT2, align=PP_ALIGN.LEFT)

# Meta info
add_textbox(slide, '2024 · 设计方案演示文稿', 1, 4.5, 8, 0.5,
            font_size=11, color=RGBColor(0x94,0xa3,0xb8), align=PP_ALIGN.LEFT)

# Decorative circle (right)
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(-0.5), Inches(3), Inches(3))
circle.fill.solid()
circle.fill.fore_color.rgb = COLOR_ACCENT
circle.line.fill.background()

# ================================================================
# Slide 2: Project Overview
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_BG)

# Title bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = COLOR_ACCENT
bar.line.fill.background()

add_textbox(slide, '项目概述', 0.5, 0.4, 9, 0.6,
            font_size=24, bold=True, color=COLOR_PRIMARY, align=PP_ALIGN.LEFT)

# Content boxes
items = [
    ('🎯', '网站定位', '技术博客 + 个人作品展示\n简洁现代（Minimal）设计风格'),
    ('📄', '核心页面', '首页（个人介绍）\n博客文章列表\n项目展示页面\n关于我页面'),
    ('🛠', '技术栈', '纯 HTML5 + CSS3 + JavaScript\nGitHub Pages 原生托管\n无框架依赖，推送即部署'),
    ('✨', '特色功能', '亮色/暗色模式切换\n完全响应式设计\n标签筛选 · 平滑过渡动画'),
]
cols = [0.3, 2.6, 4.9, 7.2]
for i, (icon, title, desc) in enumerate(items):
    col = cols[i]
    # Card
    card = add_rounded_rect(slide, col, 1.5, 2.1, 3.5, COLOR_SURFACE, COLOR_BORDER)
    # Icon
    add_textbox(slide, icon, col+0.2, 1.8, 1.7, 0.6, font_size=28, align=PP_ALIGN.LEFT)
    # Title
    add_textbox(slide, title, col+0.2, 2.5, 1.7, 0.5, font_size=14, bold=True, color=COLOR_PRIMARY)
    # Desc
    add_textbox(slide, desc, col+0.2, 3.1, 1.7, 1.5, font_size=11, color=COLOR_TEXT_LIGHT)

# ================================================================
# Slide 3: Design Style
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_BG)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = COLOR_ACCENT
bar.line.fill.background()

add_textbox(slide, '设计风格', 0.5, 0.4, 9, 0.6,
            font_size=24, bold=True, color=COLOR_PRIMARY, align=PP_ALIGN.LEFT)

# Style keywords
keywords = ['Minimal', 'Clean', 'Modern', 'Typography-focused']
x_pos = 0.5
for kw in keywords:
    label = add_rounded_rect(slide, x_pos, 1.2, len(kw)*0.35+0.5, 0.5, COLOR_ACCENT, COLOR_ACCENT)
    tf = label.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = kw
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER
    x_pos += len(kw)*0.35 + 0.8

# Color palette
add_textbox(slide, '配色方案', 0.5, 2.0, 9, 0.5, font_size=16, bold=True, color=COLOR_PRIMARY)

colors_data = [
    (COLOR_PRIMARY,  '主色\n#1a1a2e'),
    (COLOR_ACCENT,   '主题色\n#4f46e5'),
    (COLOR_ACCENT2,  '强调色\n#06b6d4'),
    (COLOR_BG,       '背景色\n#f8f9fa'),
    (COLOR_TEXT,     '文字色\n#333333'),
]
for i, (color, label) in enumerate(colors_data):
    x = 0.5 + i * 1.85
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.7), Inches(1.5), Inches(1.0))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.color.rgb = COLOR_BORDER
    rect.line.width = Pt(1)
    tf = rect.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_WHITE if color in [COLOR_PRIMARY, COLOR_ACCENT, COLOR_ACCENT2, COLOR_TEXT] else COLOR_TEXT
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(8)

# Typography
add_textbox(slide, '字体系统', 0.5, 4.0, 9, 0.5, font_size=16, bold=True, color=COLOR_PRIMARY)
add_textbox(slide, '标题：32px / Bold', 0.5, 4.5, 4, 0.4, font_size=13, color=COLOR_TEXT)
add_textbox(slide, '正文：16px / Regular', 0.5, 4.9, 4, 0.4, font_size=13, color=COLOR_TEXT_LIGHT)
add_textbox(slide, '代码：系统等宽字体', 5.5, 4.5, 4, 0.4, font_size=13, color=COLOR_TEXT)
add_textbox(slide, 'font-family: system-ui, "PingFang SC", sans-serif', 5.5, 4.9, 4.5, 0.4, font_size=10, color=COLOR_TEXT_LIGHT)

# ================================================================
# Slide 4: Page Structure
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_BG)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = COLOR_ACCENT
bar.line.fill.background()

add_textbox(slide, '页面结构与导航', 0.5, 0.4, 9, 0.6,
            font_size=24, bold=True, color=COLOR_PRIMARY, align=PP_ALIGN.LEFT)

# Site map (tree)
pages = [
    ('🌐', 'username.github.io', COLOR_PRIMARY, True),
    ('  ├─', '首页（Hero + 快捷入口 + 最新文章）', COLOR_TEXT, False),
    ('  ├─', '博客列表（文章卡片 + 标签筛选）', COLOR_TEXT, False),
    ('  ├─', '项目展示（项目卡片网格）', COLOR_TEXT, False),
    ('  └─', '关于我（个人简介 + 技能栈）', COLOR_TEXT, False),
]
y = 1.3
for prefix, name, color, is_bold in pages:
    add_textbox(slide, prefix + ' ' + name, 0.8, y, 8, 0.5,
                font_size=13, bold=is_bold, color=color)
    y += 0.6

# Navigation bar preview
nav_bar = add_rounded_rect(slide, 0.5, 3.2, 9, 0.6, RGBColor(0xff,0xff,0xff), COLOR_BORDER)
tf = nav_bar.text_frame
p = tf.paragraphs[0]
p.text = '  Logo   首页  博客  项目  关于        🌙'
p.font.size = Pt(11)
p.font.color.rgb = COLOR_TEXT
p.alignment = PP_ALIGN.LEFT

# Interaction diagram
add_textbox(slide, '交互设计', 0.5, 4.2, 9, 0.5, font_size=16, bold=True, color=COLOR_PRIMARY)
interactions = [
    '页面切换：淡入淡出过渡（0.3s）',
    '卡片 Hover：上浮 4px + 阴影加深',
    '导航链接：底部滑入下划线',
    '暗色模式：全站颜色平滑过渡（0.5s）',
]
for i, item in enumerate(interactions):
    add_textbox(slide, '• ' + item, 0.8, 4.7 + i*0.35, 8, 0.4, font_size=11, color=COLOR_TEXT_LIGHT)

# ================================================================
# Slide 5: Homepage Design
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_BG)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = COLOR_ACCENT
bar.line.fill.background()

add_textbox(slide, '首页设计', 0.5, 0.4, 9, 0.6,
            font_size=24, bold=True, color=COLOR_PRIMARY, align=PP_ALIGN.LEFT)

# Mockup frame
frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.2), Inches(7), Inches(3.8))
frame.fill.solid()
frame.fill.fore_color.rgb = COLOR_WHITE
frame.line.color.rgb = COLOR_BORDER
frame.line.width = Pt(1)

# Nav bar in mockup
nav = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.2), Inches(7), Inches(0.4))
nav.fill.solid()
nav.fill.fore_color.rgb = RGBColor(0xf1,0xf5,0xf9)
nav.line.fill.background()

# Hero section (avatar + name)
avatar = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.7), Inches(1.8), Inches(0.9), Inches(0.9))
avatar.fill.solid()
avatar.fill.fore_color.rgb = COLOR_ACCENT
avatar.line.fill.background()

add_textbox(slide, '你好，我是开发者', 2.5, 2.9, 5, 0.4, font_size=11, bold=True, color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)
add_textbox(slide, '全栈开发者 / 开源爱好者 / 技术写作', 2.5, 3.3, 5, 0.3, font_size=9, color=COLOR_TEXT_LIGHT, align=PP_ALIGN.CENTER)

# Quick link cards
for i, label in enumerate(['技术博客', '项目作品', '关于我']):
    x = 2.2 + i * 1.7
    card = add_rounded_rect(slide, x, 3.8, 1.4, 0.7, COLOR_SURFACE, COLOR_BORDER)
    tf = card.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.color.rgb = COLOR_TEXT
    p.alignment = PP_ALIGN.CENTER

# Notes
add_textbox(slide, '设计说明：', 0.5, 5.0, 9, 0.4, font_size=12, bold=True, color=COLOR_PRIMARY)
notes = [
    '• 居中布局，突出个人品牌形象',
    '• 圆形头像 + 渐变背景，增强视觉记忆',
    '• 三列快捷入口卡片，清晰引导用户',
]
for i, note in enumerate(notes):
    add_textbox(slide, note, 0.8, 5.3 + i*0.3, 9, 0.3, font_size=10, color=COLOR_TEXT_LIGHT)

# ================================================================
# Slide 6: Blog & Projects Design
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_BG)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = COLOR_ACCENT
bar.line.fill.background()

add_textbox(slide, '博客 & 项目页设计', 0.5, 0.4, 9, 0.6,
            font_size=24, bold=True, color=COLOR_PRIMARY, align=PP_ALIGN.LEFT)

# Blog section
add_textbox(slide, '📝 博客列表页', 0.5, 1.2, 4, 0.5, font_size=16, bold=True, color=COLOR_PRIMARY)

# Tag filter
for i, tag in enumerate(['全部', '前端', '后端', 'DevOps']):
    x = 0.5 + i * 0.9
    tag_btn = add_rounded_rect(slide, x, 1.8, 0.7, 0.35, COLOR_ACCENT if i==0 else COLOR_SURFACE, COLOR_ACCENT if i==0 else COLOR_BORDER)
    tf = tag_btn.text_frame
    p = tf.paragraphs[0]
    p.text = tag
    p.font.size = Pt(8)
    p.font.color.rgb = COLOR_WHITE if i==0 else COLOR_TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER

# Blog post cards
for i in range(3):
    y = 2.4 + i * 0.65
    card = add_rounded_rect(slide, 0.5, y, 4.3, 0.55, COLOR_SURFACE, COLOR_BORDER)
    tf = card.text_frame
    p = tf.paragraphs[0]
    p.text = f'文章标题示例 {i+1}'
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_TEXT
    # Date
    add_textbox(slide, '2024-03-15', 0.7, y+0.28, 1.2, 0.25, font_size=8, color=COLOR_TEXT_LIGHT)

# Projects section
add_textbox(slide, '💼 项目展示页', 5.2, 1.2, 4, 0.5, font_size=16, bold=True, color=COLOR_PRIMARY)

# Project cards (2x2 grid)
grid_positions = [(5.2, 1.8), (7.7, 1.8), (5.2, 3.5), (7.7, 3.5)]
for i, (x, y) in enumerate(grid_positions):
    card = add_rounded_rect(slide, x, y, 2.1, 1.4, COLOR_SURFACE, COLOR_BORDER)
    # Project icon area
    icon_area = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(2.1), Inches(0.6))
    icon_area.fill.solid()
    colors = [COLOR_ACCENT, COLOR_ACCENT2, RGBColor(0x16,0xa3,0x4a), RGBColor(0xd9,0x77,0x06)]
    icon_area.fill.fore_color.rgb = colors[i % len(colors)]
    icon_area.line.fill.background()
    tf = icon_area.text_frame
    p = tf.paragraphs[0]
    p.text = f'📦'
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER
    # Project name
    add_textbox(slide, f'项目 {i+1}', x+0.1, y+0.65, 1.9, 0.3, font_size=10, bold=True, color=COLOR_TEXT)
    add_textbox(slide, '项目描述文字...', x+0.1, y+0.95, 1.9, 0.3, font_size=8, color=COLOR_TEXT_LIGHT)

# ================================================================
# Slide 7: Responsive & Dark Mode
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_BG)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = COLOR_ACCENT
bar.line.fill.background()

add_textbox(slide, '响应式设计 & 暗色模式', 0.5, 0.4, 9, 0.6,
            font_size=24, bold=True, color=COLOR_PRIMARY, align=PP_ALIGN.LEFT)

# Device frames
devices = [('📱 移动端\n< 768px', 0.8), ('💻 平板\n768-1024px', 3.8), ('🖥️ 桌面端\n> 1024px', 6.8)]
for label, x in devices:
    # Device frame
    frame = add_rounded_rect(slide, x, 1.5, 2.5, 3.0, COLOR_SURFACE, COLOR_BORDER)
    tf = frame.text_frame
    p = tf.paragraphs[0]
    p.text = label.split('\n')[0]
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(4)
    # Add second line
    run = p.add_run()
    run.text = label.split('\n')[1]
    run.font.size = Pt(8)
    run.font.color.rgb = COLOR_TEXT_LIGHT
    run.font.name = 'Arial'

    # Simplified screen content
    screen = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x+0.15), Inches(2.1), Inches(2.2), Inches(2.2))
    screen.fill.solid()
    screen.fill.fore_color.rgb = RGBColor(0xf1,0xf5,0xf9)
    screen.line.color.rgb = COLOR_BORDER
    screen.line.width = Pt(0.5)

# Dark mode section
add_textbox(slide, '🌙 暗色模式', 0.5, 4.8, 9, 0.4, font_size=16, bold=True, color=COLOR_PRIMARY)

dark_colors = [
    ('背景', '#0f0f1a', RGBColor(0x0f,0x0f,0x1a)),
    ('卡片', '#1a1a2e', RGBColor(0x1a,0x1a,0x2e)),
    ('文字', '#e0e0e0', RGBColor(0xe0,0xe0,0xe0)),
]
for i, (label, hex_val, color) in enumerate(dark_colors):
    x = 0.8 + i * 3.0
    swatch = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(5.2), Inches(0.8), Inches(0.5))
    swatch.fill.solid()
    swatch.fill.fore_color.rgb = color
    swatch.line.color.rgb = COLOR_BORDER
    swatch.line.width = Pt(1)
    tf = swatch.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.color.rgb = COLOR_WHITE if color in [RGBColor(0x0f,0x0f,0x1a), RGBColor(0x1a,0x1a,0x2e)] else COLOR_TEXT
    p.alignment = PP_ALIGN.CENTER

# ================================================================
# Slide 8: Deployment Guide
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_BG)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = COLOR_ACCENT
bar.line.fill.background()

add_textbox(slide, '部署流程', 0.5, 0.4, 9, 0.6,
            font_size=24, bold=True, color=COLOR_PRIMARY, align=PP_ALIGN.LEFT)

# Steps
steps = [
    ('1', '创建 GitHub 仓库', '仓库名：username.github.io\n设置为 Public'),
    ('2', '克隆到本地', 'git clone https://github.com/\nusername/username.github.io.git'),
    ('3', '复制网站文件', '将 index.html、css/、js/ 等\n文件复制到仓库目录'),
    ('4', '自定义内容', '编辑 js/data.js 更新\n博客文章和项目数据'),
    ('5', '提交并推送', 'git add . → commit → push\n等待 1-2 分钟自动部署'),
    ('6', '访问网站', 'https://username.github.io\n（可选）配置自定义域名'),
]
for i, (num, title, desc) in enumerate(steps):
    y = 1.3 + i * 0.68
    # Step number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(y), Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLOR_ACCENT
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Title
    add_textbox(slide, title, 1.1, y, 3, 0.45, font_size=13, bold=True, color=COLOR_PRIMARY)
    # Desc
    add_textbox(slide, desc, 4.2, y-0.05, 5.3, 0.6, font_size=10, color=COLOR_TEXT_LIGHT)

# ================================================================
# Slide 9: Thank You
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_PRIMARY)

add_textbox(slide, '谢谢', 1, 1.8, 8, 1,
            font_size=44, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

add_textbox(slide, '期待你的个人网站早日上线 🚀', 1, 3.0, 8, 0.6,
            font_size=16, color=COLOR_ACCENT2, align=PP_ALIGN.CENTER)

add_textbox(slide, 'GitHub Pages + 纯静态网站 = 免费 & 高效 & 优雅', 1, 4.0, 8, 0.5,
            font_size=12, color=RGBColor(0x94,0xa3,0xb8), align=PP_ALIGN.CENTER)

# Decorative elements
for i, (x, y) in enumerate([(0.5, 0.5), (8.5, 0.5), (0.5, 4.5), (8.5, 4.5)]):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.3), Inches(0.3))
    dot.fill.solid()
    dot.fill.fore_color.rgb = COLOR_ACCENT if i % 2 == 0 else COLOR_ACCENT2
    dot.line.fill.background()

# ---- Save ----
output_path = os.path.join(os.path.dirname(__file__), '..', 'docs', 'Website Design Proposal.pptx')
output_path = os.path.abspath(output_path)
prs.save(output_path)
print('PPTX saved to: ' + output_path)
